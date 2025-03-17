import os
import re
import uuid
import logging
import numpy as np
import pandas as pd
from typing import List, Dict, Tuple, Optional, Union, Any, Iterator
from pathlib import Path
import pickle
import hashlib
import json

# Document processing libraries
import PyPDF2
from PyPDF2 import PdfReader
import docx
import pptx
import mammoth
import openpyxl
import csv
from bs4 import BeautifulSoup
import markdown
import html2text
from PIL import Image
import pytesseract
import unstructured
from unstructured.partition.auto import partition
from unstructured.partition.html import partition_html
from unstructured.partition.pdf import partition_pdf
from unstructured.partition.image import partition_image
from unstructured.cleaners.core import clean_extra_whitespace

# NLP libraries for advanced chunking
import nltk
from nltk.tokenize import sent_tokenize
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt')

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler(),
        logging.FileHandler("document_processor.log")
    ]
)
logger = logging.getLogger("DocumentProcessor")

class DocumentProcessor:
    """
    Processes various document types, handles chunking with context preservation,
    and prepares content for embedding.
    """
    
    def __init__(self, 
                 embedding_model_path: str,
                 chunk_size: int = 1000,
                 chunk_overlap: int = 200,
                 min_chunk_size: int = 100,
                 extract_tables: bool = True,
                 extract_images: bool = True):
        """
        Initialize the document processor.
        
        Args:
            embedding_model_path: Path to the embedding model
            chunk_size: Target size for text chunks
            chunk_overlap: Number of characters to overlap between chunks
            min_chunk_size: Minimum chunk size to consider valid
            extract_tables: Whether to extract tables from documents
            extract_images: Whether to extract and OCR images
        """
        self.embedding_model_path = embedding_model_path
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.min_chunk_size = min_chunk_size
        self.extract_tables = extract_tables
        self.extract_images = extract_images
        self._load_embedding_model()
        
    def _load_embedding_model(self):
        """Load the embedding model."""
        try:
            from sentence_transformers import SentenceTransformer
            
            if os.path.exists(self.embedding_model_path):
                logger.info(f"Loading embedding model from {self.embedding_model_path}")
                self.embedding_model = SentenceTransformer(self.embedding_model_path)
            else:
                # Fall back to a default model
                logger.warning(f"Model not found at {self.embedding_model_path}, using default model")
                self.embedding_model = SentenceTransformer("all-MiniLM-L6-v2")
            
            # Get embedding dimension
            self.embedding_dim = self.embedding_model.get_sentence_embedding_dimension()
            logger.info(f"Embedding model loaded with dimension {self.embedding_dim}")
        except Exception as e:
            logger.error(f"Error loading embedding model: {e}")
            raise
    
    def process_file(self, file_path: Union[str, Path]) -> Dict:
        """
        Process a file and extract content with metadata.
        
        Args:
            file_path: Path to the file
        
        Returns:
            Dict with file ID, content, and metadata
        """
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        # Generate a unique file ID based on path and modification time
        file_stats = file_path.stat()
        file_id = hashlib.md5(f"{file_path}_{file_stats.st_mtime}".encode()).hexdigest()
        
        # Determine file type and extract content
        file_type = file_path.suffix.lower()
        
        try:
            if file_type in ['.pdf']:
                content, metadata = self._process_pdf(file_path)
            elif file_type in ['.docx', '.doc']:
                content, metadata = self._process_word(file_path)
            elif file_type in ['.pptx', '.ppt']:
                content, metadata = self._process_powerpoint(file_path)
            elif file_type in ['.xlsx', '.xls', '.csv']:
                content, metadata = self._process_spreadsheet(file_path)
            elif file_type in ['.txt', '.md', '.rst']:
                content, metadata = self._process_text(file_path)
            elif file_type in ['.html', '.htm']:
                content, metadata = self._process_html(file_path)
            elif file_type in ['.jpg', '.jpeg', '.png', '.tiff', '.bmp', '.gif']:
                content, metadata = self._process_image(file_path)
            elif file_type == '.eod':  # Example of handling custom format
                content, metadata = self._process_eod(file_path)
            else:
                # Try using unstructured as fallback
                try:
                    elements = partition(filename=str(file_path))
                    content = "\n\n".join([clean_extra_whitespace(el.text) for el in elements])
                    metadata = {"title": file_path.stem, "format": "unknown"}
                except Exception as e:
                    logger.warning(f"Unstructured fallback failed: {e}")
                    content = f"Unsupported file type: {file_type}"
                    metadata = {"title": file_path.stem, "format": "unsupported"}
            
            # Add file metadata
            metadata.update({
                "filename": file_path.name,
                "file_path": str(file_path),
                "file_type": file_type,
                "file_size": file_stats.st_size,
                "modified_date": file_stats.st_mtime
            })
            
            logger.info(f"Processed {file_path.name} ({file_type}) with {len(content)} characters")
            
            return {
                "file_id": file_id,
                "content": content,
                "metadata": metadata
            }
        
        except Exception as e:
            logger.error(f"Error processing file {file_path}: {e}")
            raise
    
    def _process_pdf(self, file_path: Path) -> Tuple[str, Dict]:
        """Process PDF documents."""
        try:
            # First try with unstructured for better structure preservation
            elements = partition_pdf(filename=str(file_path), 
                                    extract_images=self.extract_images,
                                    extract_tables=self.extract_tables)
            
            # Extract titles, headers, paragraphs, tables separately
            titles = []
            headings = []
            paragraphs = []
            tables = []
            image_captions = []
            
            for element in elements:
                element_text = clean_extra_whitespace(element.text)
                if hasattr(element, 'category'):
                    if element.category == 'Title':
                        titles.append(element_text)
                    elif element.category == 'NarrativeText':
                        paragraphs.append(element_text)
                    elif element.category == 'ListItem':
                        paragraphs.append(f"• {element_text}")
                    elif element.category == 'Header':
                        headings.append(element_text)
                    elif element.category == 'Table':
                        tables.append(element_text)
                    elif element.category == 'Image':
                        if hasattr(element, 'caption') and element.caption:
                            image_captions.append(f"Image: {element.caption}")
                        else:
                            image_captions.append(f"Image on page {element.metadata.page_number if hasattr(element.metadata, 'page_number') else 'unknown'}")
                else:
                    paragraphs.append(element_text)
            
            # Fallback to PyPDF2 if unstructured didn't extract much
            if not paragraphs and not tables:
                reader = PdfReader(file_path)
                metadata = reader.metadata or {}
                content = ""
                for i, page in enumerate(reader.pages):
                    page_text = page.extract_text() or ""
                    if page_text.strip():
                        content += f"\n\n--- Page {i+1} ---\n\n{page_text}"
                
                # Try to detect and preserve structural elements
                lines = content.split('\n')
                for i, line in enumerate(lines):
                    line = line.strip()
                    if not line:
                        continue
                    
                    # Heuristically detect headings and titles
                    if len(line) < 100 and line.endswith(':'):
                        headings.append(line)
                    elif len(line) < 80 and line.isupper():
                        titles.append(line)
                    elif len(line) < 120 and i > 0 and not lines[i-1].strip():
                        # Likely a heading or title if preceded by empty line and short
                        headings.append(line)
                    else:
                        paragraphs.append(line)
            
            # Combine all elements with structure markers
            structured_content = ""
            
            if titles:
                structured_content += "# TITLES\n" + "\n".join(titles) + "\n\n"
            
            if headings:
                structured_content += "# HEADINGS\n" + "\n".join(headings) + "\n\n"
            
            if paragraphs:
                structured_content += "# CONTENT\n" + "\n\n".join(paragraphs) + "\n\n"
            
            if tables:
                structured_content += "# TABLES\n" + "\n\n".join(tables) + "\n\n"
            
            if image_captions:
                structured_content += "# IMAGES\n" + "\n".join(image_captions) + "\n\n"
            
            # Extract PDF metadata
            pdf_metadata = {}
            try:
                reader = PdfReader(file_path)
                info = reader.metadata
                if info:
                    pdf_metadata = {
                        "title": info.title,
                        "author": info.author,
                        "subject": info.subject,
                        "creator": info.creator,
                        "producer": info.producer,
                        "page_count": len(reader.pages)
                    }
                    # Clean out None values
                    pdf_metadata = {k: v for k, v in pdf_metadata.items() if v is not None}
            except Exception as e:
                logger.warning(f"Error extracting PDF metadata: {e}")
            
            return structured_content, pdf_metadata
            
        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {e}")
            # Fallback to basic extraction
            content = ""
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    content += page.extract_text() + "\n\n"
            return content, {"title": file_path.stem, "format": "pdf"}
    
    def _process_word(self, file_path: Path) -> Tuple[str, Dict]:
        """Process Word documents."""
        try:
            # Try mammoth first for better HTML conversion
            with open(file_path, "rb") as docx_file:
                result = mammoth.convert_to_html(docx_file)
                html = result.value
                
                # Convert HTML to markdown-like structure with preserving headers
                h = html2text.HTML2Text()
                h.ignore_links = False
                h.ignore_images = False
                h.ignore_tables = False
                h.body_width = 0  # No wrapping
                structured_text = h.handle(html)
                
                # Extract metadata using python-docx
                doc = docx.Document(file_path)
                metadata = {
                    "title": file_path.stem,
                    "format": "docx",
                    "paragraph_count": len(doc.paragraphs),
                    "section_count": len(doc.sections)
                }
                
                # Try to extract core properties
                try:
                    core_props = doc.core_properties
                    metadata.update({
                        "author": core_props.author,
                        "created": core_props.created,
                        "modified": core_props.modified,
                        "title": core_props.title or file_path.stem
                    })
                    # Clean out None values
                    metadata = {k: v for k, v in metadata.items() if v is not None}
                except Exception as e:
                    logger.warning(f"Error extracting DOCX metadata: {e}")
                
                return structured_text, metadata
                
        except Exception as e:
            logger.error(f"Error processing Word doc {file_path}: {e}")
            # Fallback to python-docx
            try:
                doc = docx.Document(file_path)
                content = "\n\n".join([para.text for para in doc.paragraphs])
                return content, {"title": file_path.stem, "format": "docx"}
            except Exception as e2:
                logger.error(f"Fallback processing of Word doc failed: {e2}")
                return f"Error processing {file_path.name}: {e}", {"title": file_path.stem, "format": "docx"}
    
    def _process_powerpoint(self, file_path: Path) -> Tuple[str, Dict]:
        """Process PowerPoint presentations."""
        try:
            presentation = pptx.Presentation(file_path)
            content = []
            
            for i, slide in enumerate(presentation.slides):
                slide_content = [f"--- Slide {i+1} ---"]
                
                # Extract slide title
                if slide.shapes.title:
                    slide_content.append(f"# {slide.shapes.title.text}")
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        # Skip if it's the title we already added
                        if shape == slide.shapes.title:
                            continue
                        text = shape.text.strip()
                        if text:
                            slide_content.append(text)
                
                content.append("\n".join(slide_content))
            
            metadata = {
                "title": file_path.stem,
                "format": "pptx",
                "slide_count": len(presentation.slides)
            }
            
            return "\n\n".join(content), metadata
            
        except Exception as e:
            logger.error(f"Error processing PowerPoint {file_path}: {e}")
            return f"Error processing {file_path.name}: {e}", {"title": file_path.stem, "format": "pptx"}
    
    def _process_spreadsheet(self, file_path: Path) -> Tuple[str, Dict]:
        """Process Excel and CSV files."""
        try:
            file_type = file_path.suffix.lower()
            
            if file_type in ['.csv']:
                # Process CSV
                tables = []
                
                # Read CSV with pandas to handle various formats and encodings
                try:
                    df = pd.read_csv(file_path, encoding='utf-8')
                except UnicodeDecodeError:
                    # Try different encodings if utf-8 fails
                    df = pd.read_csv(file_path, encoding='latin1')
                
                # Convert each table to text with headers
                tables.append(f"Table from {file_path.name}:\n{df.to_string(index=False)}")
                
                # Create metadata about the table
                table_metadata = {
                    "title": file_path.stem,
                    "format": "csv",
                    "columns": df.columns.tolist(),
                    "row_count": len(df),
                    "column_count": len(df.columns)
                }
                
                return "\n\n".join(tables), table_metadata
                
            elif file_type in ['.xlsx', '.xls']:
                # Process Excel
                try:
                    # Read all sheets
                    xlsx = pd.ExcelFile(file_path)
                    sheet_names = xlsx.sheet_names
                    
                    tables = []
                    for sheet_name in sheet_names:
                        df = pd.read_excel(xlsx, sheet_name)
                        if not df.empty:
                            tables.append(f"Sheet: {sheet_name}\n{df.to_string(index=False)}")
                    
                    # Get basic metadata
                    excel_metadata = {
                        "title": file_path.stem,
                        "format": "excel",
                        "sheets": sheet_names,
                        "sheet_count": len(sheet_names)
                    }
                    
                    # Try to get more detailed Excel metadata
                    try:
                        workbook = openpyxl.load_workbook(file_path, read_only=True)
                        properties = workbook.properties
                        if properties:
                            excel_metadata.update({
                                "author": properties.creator,
                                "title": properties.title or file_path.stem,
                                "created": properties.created,
                                "modified": properties.modified
                            })
                    except Exception as e:
                        logger.warning(f"Error extracting Excel metadata: {e}")
                    
                    # Clean out None values
                    excel_metadata = {k: v for k, v in excel_metadata.items() if v is not None}
                    
                    return "\n\n".join(tables), excel_metadata
                    
                except Exception as e:
                    logger.error(f"Error processing Excel file {file_path}: {e}")
                    return f"Error processing {file_path.name}: {e}", {"title": file_path.stem, "format": "excel"}
            
            else:
                return f"Unsupported spreadsheet type: {file_type}", {"title": file_path.stem, "format": "unknown"}
                
        except Exception as e:
            logger.error(f"Error processing spreadsheet {file_path}: {e}")
            return f"Error processing {file_path.name}: {e}", {"title": file_path.stem, "format": "spreadsheet"}
    
    def _process_text(self, file_path: Path) -> Tuple[str, Dict]:
        """Process plain text and markdown files."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            file_type = file_path.suffix.lower()
            if file_type == '.md':
                # Convert markdown to text with structure preservation
                html = markdown.markdown(content)
                h = html2text.HTML2Text()
                h.body_width = 0  # No wrapping
                content = h.handle(html)
                
            metadata = {
                "title": file_path.stem,
                "format": "markdown" if file_type == '.md' else "text",
                "char_count": len(content),
                "line_count": len(content.splitlines())
            }
            
            return content, metadata
            
        except UnicodeDecodeError:
            # Try different encoding
            try:
                with open(file_path, 'r', encoding='latin1') as f:
                    content = f.read()
                metadata = {
                    "title": file_path.stem,
                    "format": "text",
                    "encoding": "latin1"
                }
                return content, metadata
            except Exception as e:
                logger.error(f"Error processing text file {file_path}: {e}")
                return f"Error processing {file_path.name}: {e}", {"title": file_path.stem, "format": "text"}
    
    def _process_html(self, file_path: Path) -> Tuple[str, Dict]:
        """Process HTML files."""
        try:
            # First try with unstructured for better tag handling
            elements = partition_html(filename=str(file_path))
            
            if elements:
                # Categorize elements
                titles = []
                headings = []
                paragraphs = []
                lists = []
                tables = []
                
                for element in elements:
                    element_text = clean_extra_whitespace(element.text)
                    if hasattr(element, 'category'):
                        if element.category == 'Title':
                            titles.append(element_text)
                        elif element.category == 'NarrativeText':
                            paragraphs.append(element_text)
                        elif element.category == 'ListItem':
                            lists.append(f"• {element_text}")
                        elif element.category == 'Header':
                            headings.append(element_text)
                        elif element.category == 'Table':
                            tables.append(element_text)
                        else:
                            paragraphs.append(element_text)
                    else:
                        paragraphs.append(element_text)
                
                # Build structured content
                structured_content = ""
                
                if titles:
                    structured_content += "# TITLES\n" + "\n".join(titles) + "\n\n"
                
                if headings:
                    structured_content += "# HEADINGS\n" + "\n".join(headings) + "\n\n"
                
                if lists:
                    structured_content += "# LISTS\n" + "\n".join(lists) + "\n\n"
                
                if paragraphs:
                    structured_content += "# CONTENT\n" + "\n\n".join(paragraphs) + "\n\n"
                
                if tables:
                    structured_content += "# TABLES\n" + "\n\n".join(tables) + "\n\n"
                
                metadata = {
                    "title": file_path.stem,
                    "format": "html"
                }
                
                return structured_content, metadata
            
            # Fallback to BeautifulSoup
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            soup = BeautifulSoup(content, 'html.parser')
            
            # Extract title
            title = soup.title.string if soup.title else file_path.stem
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.extract()
            
            # Get text
            text = soup.get_text(separator='\n')
            
            # Format text
            lines = [line.strip() for line in text.splitlines() if line.strip()]
            
            metadata = {
                "title": title,
                "format": "html",
                "headings": [h.text for h in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6'])]
            }
            
            return "\n\n".join(lines), metadata
            
        except Exception as e:
            logger.error(f"Error processing HTML file {file_path}: {e}")
            return f"Error processing {file_path.name}: {e}", {"title": file_path.stem, "format": "html"}
    
    def _process_image(self, file_path: Path) -> Tuple[str, Dict]:
        """Process image files with OCR."""
        try:
            if not self.extract_images:
                return f"Image file: {file_path.name} (OCR disabled)", {"title": file_path.stem, "format": "image"}
            
            # First try with unstructured
            try:
                elements = partition_image(filename=str(file_path))
                if elements:
                    text = "\n\n".join([clean_extra_whitespace(el.text) for el in elements if hasattr(el, 'text') and el.text])
                    if text.strip():
                        metadata = {
                            "title": file_path.stem,
                            "format": "image",
                            "extraction_method": "unstructured"
                        }
                        return text, metadata
            except Exception as e:
                logger.warning(f"Unstructured image processing failed: {e}, falling back to pytesseract")
            
            # Fallback to pytesseract
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            
            metadata = {
                "title": file_path.stem,
                "format": "image",
                "image_size": f"{image.width}x{image.height}",
                "extraction_method": "pytesseract"
            }
            
            return text, metadata
            
        except Exception as e:
            logger.error(f"Error processing image file {file_path}: {e}")
            return f"Error processing image {file_path.name}: {e}", {"title": file_path.stem, "format": "image"}
    
    def _process_eod(self, file_path: Path) -> Tuple[str, Dict]:
        """Process custom EOD format files."""
        try:
            # Assuming EOD is some custom format - example implementation
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # Process according to EOD format specifications
            # This is a placeholder - replace with actual EOD processing
            sections = re.split(r'={3,}', content)
            processed_sections = []
            
            for section in sections:
                section = section.strip()
                if section:
                    processed_sections.append(section)
            
            metadata = {
                "title": file_path.stem,
                "format": "eod",
                "section_count": len(processed_sections)
            }
            
            return "\n\n".join(processed_sections), metadata
            
        except Exception as e:
            logger.error(f"Error processing EOD file {file_path}: {e}")
            return f"Error processing {file_path.name}: {e}", {"title": file_path.stem, "format": "eod"}
    
    def chunk_document(self, 
                     document: Dict, 
                     chunk_by_title: bool = True) -> List[Dict]:
        """
        Chunk document content with context preservation.
        
        Args:
            document: Document dict with content and metadata
            chunk_by_title: Whether to use titles/headers for chunk boundaries
        
        Returns:
            List of chunk dicts with text and metadata
        """
        try:
            content = document["content"]
            metadata = document["metadata"]
            file_id = document["file_id"]
            
            chunks = []
            
            # Special handling for different document types
            doc_type = metadata.get("format", "")
            
            if doc_type in ["csv", "excel"]:
                # For tabular data, chunk by table/sheet
                chunks = self._chunk_tabular(content, metadata, file_id)
            elif doc_type == "pdf" and "TABLES" in content:
                # Special handling for PDFs with detected tables
                chunks = self._chunk_pdf_with_tables(content, metadata, file_id)
            elif chunk_by_title:
                # Chunk by sections using titles and headings
                chunks = self._chunk_by_section(content, metadata, file_id)
            else:
                # Default chunking method
                chunks = self._chunk_text(content, metadata, file_id)
            
            # Ensure all chunks have required metadata
            for chunk in chunks:
                chunk["metadata"]["file_id"] = file_id
                if "chunk_id" not in chunk["metadata"]:
                    chunk["metadata"]["chunk_id"] = str(uuid.uuid4())
            
            logger.info(f"Created {len(chunks)} chunks from document {metadata.get('filename', 'unknown')}")
            return chunks
            
        except Exception as e:
            logger.error(f"Error chunking document: {e}")
            # Create a single chunk with error information
            return [{
                "text": f"Error chunking document: {e}",
                "metadata": {
                    "file_id": document.get("file_id", "unknown"),
                    "chunk_id": str(uuid.uuid4()),
                    "error": str(e)
                }
            }]
    
    def _chunk_tabular(self, content: str, metadata: Dict, file_id: str) -> List[Dict]:
        """Chunk tabular data from CSV/Excel."""
        chunks = []
        
        # Split by tables/sheets if multiple are present
        tables = re.split(r'(?:Sheet:|Table from)', content)
        
        for i, table in enumerate(tables):
            table = table.strip()
            if not table:
                continue
                
            # Add back the "Table from" or "Sheet:" prefix if it was removed
            if i > 0:
                # Check if this looks like a sheet name
                if re.match(r'^\w+\s*\n', table):
                    sheet_name = table.split('\n')[0].strip()
                    table = f"Sheet: {table}"
                else:
                    table = f"Table from {table}"
            
            # Check if the table is too large
            if len(table) > self.chunk_size * 2:
                # Split large tables by rows
                rows = table.split('\n')
                header = rows[0] if rows else ""
                
                # Group rows into chunks
                current_chunk = [header]
                current_length = len(header)
                
                for row in rows[1:]:
                    # Always include header with each chunk
                    if current_length + len(row) > self.chunk_size and len(current_chunk) > 1:
                        chunk_text = '\n'.join(current_chunk)
                        chunks.append({
                            "text": chunk_text,
                            "metadata": {
                                "chunk_id": str(uuid.uuid4()),
                                "file_id": file_id,
                                "chunk_type": "table",
                                "table_index": i,
                                "row_count": len(current_chunk) - 1,  # Subtract header
                                "is_partial": True,
                                **{k: v for k, v in metadata.items() if k not in ["chunk_id", "file_id"]}
                            }
                        })
                        # Start new chunk with header
                        current_chunk = [header, row]
                        current_length = len(header) + len(row)
                    else:
                        current_chunk.append(row)
                        current_length += len(row)
                
                # Add the last chunk
                if len(current_chunk) > 1:
                    chunk_text = '\n'.join(current_chunk)
                    chunks.append({
                        "text": chunk_text,
                        "metadata": {
                            "chunk_id": str(uuid.uuid4()),
                            "file_id": file_id,
                            "chunk_type": "table",
                            "table_index": i,
                            "row_count": len(current_chunk) - 1,  # Subtract header
                            "is_partial": True,
                            **{k: v for k, v in metadata.items() if k not in ["chunk_id", "file_id"]}
                        }
                    })
            else:
                # Table fits in a single chunk
                chunks.append({
                    "text": table,
                    "metadata": {
                        "chunk_id": str(uuid.uuid4()),
                        "file_id": file_id,
                        "chunk_type": "table",
                        "table_index": i,
                        **{k: v for k, v in metadata.items() if k not in ["chunk_id", "file_id"]}
                    }
                })
        
        return chunks
    
    def _chunk_pdf_with_tables(self, content: str, metadata: Dict, file_id: str) -> List[Dict]:
        """Chunk PDF content that contains tables."""
        chunks = []
        
        # Split by sections
        sections = re.split(r'# (TITLES|HEADINGS|CONTENT|TABLES|IMAGES)', content)
        
        current_section_type = None
        current_text = ""
        
        # Process each section
        for i, section in enumerate(sections):
            if section in ["TITLES", "HEADINGS", "CONTENT", "TABLES", "IMAGES"]:
                current_section_type = section
                continue
            
            if not current_section_type or not section.strip():
                continue
            
            # Handle each section based on type
            if current_section_type == "TABLES":
                # Process tables separately
                tables = re.split(r'\n\n+', section)
                for table in tables:
                    if not table.strip():
                        continue
                    
                    chunks.append({
                        "text": table,
                        "metadata": {
                            "chunk_id": str(uuid.uuid4()),
                            "file_id": file_id,
                            "chunk_type": "table",
                            **{k: v for k, v in metadata.items() if k not in ["chunk_id", "file_id"]}
                        }
                    })
            elif current_section_type in ["TITLES", "HEADINGS"]:
                # Collect titles and headings to add as context to content chunks
                current_text += section + "\n\n"
            elif current_section_type == "IMAGES":
                # Add image descriptions as separate chunks
                images = re.split(r'\n+', section)
                for image in images:
                    if not image.strip():
                        continue
                    
                    chunks.append({
                        "text": image,
                        "metadata": {
                            "chunk_id": str(uuid.uuid4()),
                            "file_id": file_id,
                            "chunk_type": "image",
                            **{k: v for k, v in metadata.items() if k not in ["chunk_id", "file_id"]}
                        }
                    })
            else:  # CONTENT
                # Add the actual content to current_text
                current_text += section
        
        # Chunk the non-table content
        if current_text:
            content_chunks = self._chunk_text(current_text, metadata, file_id)
            chunks.extend(content_chunks)
        
        return chunks
    
    def _chunk_by_section(self, content: str, metadata: Dict, file_id: str) -> List[Dict]:
        """Chunk content by natural sections using headings and structure."""
        chunks = []
        
        # Look for headings in text (common patterns)
        heading_patterns = [
            r'^#{1,6} .+,  # Markdown headings
            r'^[A-Z][A-Za-z0-9 ]+:,  # Title with colon
            r'^[0-9]+\.[0-9.]* .+,  # Numbered sections like "1.2.3 Title"
            r'^(?:Section|Chapter|Part) [0-9]+:? .+,  # Named sections
            r'^[A-Z][A-Z ]+  # ALL CAPS titles
        ]
        
        # Check if content has structured sections marked with # TITLES, # HEADINGS etc.
        if "# TITLES" in content or "# HEADINGS" in content:
            return self._chunk_pdf_with_tables(content, metadata, file_id)
        
        # Try to split by headings first
        combined_pattern = '|'.join(f'({p})' for p in heading_patterns)
        matches = list(re.finditer(fr'(?:^|\n)({combined_pattern})(?:\n|$)', content, re.MULTILINE))
        
        if matches and len(matches) > 1:
            # We found headings to use as section boundaries
            sections = []
            
            # Get the text before the first heading
            if matches[0].start() > 0:
                sections.append({
                    "heading": "Introduction",
                    "text": content[:matches[0].start()].strip()
                })
            
            # Process each heading and its content
            for i, match in enumerate(matches):
                heading = match.group(0).strip()
                
                # Get the section text - from after this heading to the next one (or end)
                start = match.end()
                end = matches[i+1].start() if i < len(matches)-1 else len(content)
                
                section_text = content[start:end].strip()
                if heading and section_text:
                    sections.append({
                        "heading": heading,
                        "text": section_text
                    })
            
            # Create chunks from sections
            for section in sections:
                section_text = f"{section['heading']}\n\n{section['text']}"
                
                if len(section_text) <= self.chunk_size:
                    # Section fits in one chunk
                    chunks.append({
                        "text": section_text,
                        "metadata": {
                            "chunk_id": str(uuid.uuid4()),
                            "file_id": file_id,
                            "chunk_type": "section",
                            "heading": section['heading'],
                            **{k: v for k, v in metadata.items() if k not in ["chunk_id", "file_id"]}
                        }
                    })
                else:
                    # Section needs further chunking
                    sub_chunks = self._chunk_text(section_text, metadata, file_id)
                    
                    # Add the heading metadata to all sub-chunks
                    for chunk in sub_chunks:
                        chunk["metadata"]["heading"] = section['heading']
                        chunk["metadata"]["chunk_type"] = "section"
                    
                    chunks.extend(sub_chunks)
            
            return chunks
        
        # Fallback to normal chunking if no clear sections found
        return self._chunk_text(content, metadata, file_id)
    
    def _chunk_text(self, content: str, metadata: Dict, file_id: str) -> List[Dict]:
        """Chunk content by size with sentence boundary preservation."""
        chunks = []
        
        # Ensure we're working with clean text
        content = content.strip()
        
        if not content:
            return []
        
        # If the content fits into a single chunk, don't split it
        if len(content) <= self.chunk_size:
            chunks.append({
                "text": content,
                "metadata": {
                    "chunk_id": str(uuid.uuid4()),
                    "file_id": file_id,
                    "chunk_type": "text",
                    **{k: v for k, v in metadata.items() if k not in ["chunk_id", "file_id"]}
                }
            })
            return chunks
        
        # Split text into sentences
        sentences = sent_tokenize(content)
        
        # Group sentences into chunks
        current_chunk = []
        current_chunk_size = 0
        
        for sentence in sentences:
            sentence = sentence.strip()
            sentence_size = len(sentence)
            
            # If a single sentence exceeds the chunk size, split it by punctuation or words
            if sentence_size > self.chunk_size:
                # If the current chunk has content, add it first
                if current_chunk:
                    chunk_text = ' '.join(current_chunk)
                    chunks.append({
                        "text": chunk_text,
                        "metadata": {
                            "chunk_id": str(uuid.uuid4()),
                            "file_id": file_id,
                            "chunk_type": "text",
                            **{k: v for k, v in metadata.items() if k not in ["chunk_id", "file_id"]}
                        }
                    })
                    current_chunk = []
                    current_chunk_size = 0
                
                # Split the long sentence
                split_parts = re.split(r'[,;:)\]}"]\s+', sentence)
                current_part = []
                current_part_size = 0
                
                for part in split_parts:
                    if current_part_size + len(part) + 1 > self.chunk_size and current_part:
                        # This part would make the chunk too big, save current and start new
                        part_text = ' '.join(current_part)
                        chunks.append({
                            "text": part_text,
                            "metadata": {
                                "chunk_id": str(uuid.uuid4()),
                                "file_id": file_id,
                                "chunk_type": "text",
                                "is_partial_sentence": True,
                                **{k: v for k, v in metadata.items() if k not in ["chunk_id", "file_id"]}
                            }
                        })
                        current_part = [part]
                        current_part_size = len(part)
                    else:
                        current_part.append(part)
                        current_part_size += len(part) + 1  # +1 for the space
                
                # Don't forget the last part
                if current_part:
                    part_text = ' '.join(current_part)
                    chunks.append({
                        "text": part_text,
                        "metadata": {
                            "chunk_id": str(uuid.uuid4()),
                            "file_id": file_id,
                            "chunk_type": "text",
                            "is_partial_sentence": True,
                            **{k: v for k, v in metadata.items() if k not in ["chunk_id", "file_id"]}
                        }
                    })
            
            # Normal case: sentence fits in a chunk
            elif current_chunk_size + sentence_size + 1 <= self.chunk_size:
                # This sentence fits in the current chunk
                current_chunk.append(sentence)
                current_chunk_size += sentence_size + 1  # +1 for the space
            
            # This sentence would make the chunk too big
            elif current_chunk_size + sentence_size + 1 > self.chunk_size:
                # If adding to current chunk exceeds size, start a new chunk
                
                # But first save the current chunk
                if current_chunk:
                    chunk_text = ' '.join(current_chunk)
                    chunks.append({
                        "text": chunk_text,
                        "metadata": {
                            "chunk_id": str(uuid.uuid4()),
                            "file_id": file_id,
                            "chunk_type": "text",
                            **{k: v for k, v in metadata.items() if k not in ["chunk_id", "file_id"]}
                        }
                    })
                
                # Add overlap from the end of the previous chunk if possible
                if current_chunk and self.chunk_overlap > 0:
                    overlap_text = ' '.join(current_chunk[-3:]) if len(current_chunk) > 3 else ' '.join(current_chunk)
                    current_chunk = [overlap_text, sentence]
                    current_chunk_size = len(overlap_text) + sentence_size + 1
                else:
                    # Start fresh with the current sentence
                    current_chunk = [sentence]
                    current_chunk_size = sentence_size
        
        # Don't forget the last chunk
        if current_chunk:
            chunk_text = ' '.join(current_chunk)
            chunks.append({
                "text": chunk_text,
                "metadata": {
                    "chunk_id": str(uuid.uuid4()),
                    "file_id": file_id,
                    "chunk_type": "text",
                    **{k: v for k, v in metadata.items() if k not in ["chunk_id", "file_id"]}
                }
            })
        
        return chunks
    
    def embed_chunks(self, chunks: List[Dict]) -> List[Dict]:
        """
        Generate embeddings for a list of chunks.
        
        Args:
            chunks: List of chunk dicts with text and metadata
        
        Returns:
            List of chunks with embeddings added
        """
        try:
            texts = [chunk["text"] for chunk in chunks]
            
            # Generate embeddings in batches to avoid memory issues
            batch_size = 16
            all_embeddings = []
            
            for i in range(0, len(texts), batch_size):
                batch_texts = texts[i:i+batch_size]
                batch_embeddings = self.embedding_model.encode(batch_texts)
                all_embeddings.extend(batch_embeddings)
            
            # Add embeddings to chunks
            for i, embedding in enumerate(all_embeddings):
                chunks[i]["embedding"] = embedding
            
            logger.info(f"Generated embeddings for {len(chunks)} chunks")
            return chunks
            
        except Exception as e:
            logger.error(f"Error generating embeddings: {e}")
            # Return chunks without embeddings
            for chunk in chunks:
                chunk["embedding_error"] = str(e)
            return chunks
    
    def process_and_embed(self, file_path: Union[str, Path], chunk_by_title: bool = True) -> Dict:
        """
        Process a file, chunk it, and generate embeddings.
        
        Args:
            file_path: Path to the file
            chunk_by_title: Whether to use titles/headers for chunk boundaries
            
        Returns:
            Dict with file ID, chunks, and metadata
        """
        try:
            # Process the file
            document = self.process_file(file_path)
            
            # Chunk the document
            chunks = self.chunk_document(document, chunk_by_title=chunk_by_title)
            
            # Generate embeddings
            chunks_with_embeddings = self.embed_chunks(chunks)
            
            return {
                "file_id": document["file_id"],
                "chunks": chunks_with_embeddings,
                "metadata": document["metadata"]
            }
            
        except Exception as e:
            logger.error(f"Error in process_and_embed for {file_path}: {e}")
            raise
